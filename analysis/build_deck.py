"""Build refactored executive summary deck — Troy SD K-5 ELA.
Incorporates SEDA 2025.1 (post-COVID), subgroup analysis, refreshed recommendation.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Brand
TROY_BLUE = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT_RED = RGBColor(0xB0, 0x21, 0x21)
ACCENT_GREEN = RGBColor(0x1F, 0x7A, 0x3D)
ACCENT_ORANGE = RGBColor(0xCC, 0x6A, 0x11)
GRAY_DARK = RGBColor(0x33, 0x33, 0x33)
GRAY_MID = RGBColor(0x70, 0x70, 0x70)
GRAY_LIGHT = RGBColor(0xEE, 0xEE, 0xEE)
LIGHT_RED = RGBColor(0xFC, 0xE4, 0xE4)
LIGHT_GREEN = RGBColor(0xE6, 0xF4, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ---------- helpers ----------
def add_rect(slide, left, top, width, height, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s

def add_text(slide, left, top, width, height, text, size=14, bold=False,
             color=GRAY_DARK, align=PP_ALIGN.LEFT, font="Calibri", italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
    return tb

def title_bar(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.85), TROY_BLUE)
    add_text(slide, Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.5),
             title, size=22, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.4), Inches(0.52), Inches(12.5), Inches(0.32),
                 subtitle, size=11, color=RGBColor(0xCC, 0xDD, 0xEE))

def footer(slide, page_num, total=29):
    add_text(slide, Inches(0.4), Inches(7.22), Inches(8), Inches(0.25),
             "Troy SD K-5 ELA — Executive Summary  •  github.com/akarpo/tsd-k5ela-choice",
             size=8, color=GRAY_MID)
    add_text(slide, Inches(11.5), Inches(7.22), Inches(1.5), Inches(0.25),
             f"{page_num} / {total}", size=8, color=GRAY_MID, align=PP_ALIGN.RIGHT)

def add_pic(slide, path, left, top, width=None, height=None):
    return slide.shapes.add_picture(path, left, top, width=width, height=height)

import os as _os
CHART_DIR = _os.environ.get('CHART_DIR', _os.path.join(_os.path.dirname(_os.path.abspath(__file__)), '..', 'charts'))

# =================================================================
# SLIDE 1 — TITLE
# =================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), TROY_BLUE)
add_text(s, Inches(0.8), Inches(2.2), Inches(11.7), Inches(0.9),
         "Troy SD K-5 ELA", size=44, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.1), Inches(11.7), Inches(0.7),
         "The Quantitative Case for Curriculum Change", size=28, color=RGBColor(0xCC,0xDD,0xEE))
add_text(s, Inches(0.8), Inches(3.95), Inches(11.7), Inches(0.5),
         "Pressure-testing Collaborative Literacy + UFLI Foundations", size=18, color=RGBColor(0xAA,0xBB,0xCC))
add_rect(s, Inches(0.8), Inches(5.4), Inches(2.5), Inches(0.04), WHITE)
add_text(s, Inches(0.8), Inches(5.55), Inches(11.7), Inches(0.4),
         "Executive Summary — Citations + Appendix", size=15, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(5.95), Inches(11.7), Inches(0.3),
         "Prepared May 2026  •  50-district national benchmark  •  2009-2025 coverage",
         size=11, color=RGBColor(0xAA,0xBB,0xCC))
add_text(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.3),
         "Repo: github.com/akarpo/tsd-k5ela-choice    Live: tsd-k5ela-choice.karpowitsch.org",
         size=10, color=RGBColor(0x88,0xAA,0xCC), italic=True)

# =================================================================
# SLIDE 2 — THE QUESTION & ANSWER
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "The question — and the one-line answer",
          "Three independent pieces of evidence converge on the same answer")

# Question block (compact)
add_text(s, Inches(0.5), Inches(1.0), Inches(12), Inches(0.32),
         "QUESTION", size=10, bold=True, color=GRAY_MID)
add_text(s, Inches(0.5), Inches(1.32), Inches(12.3), Inches(0.7),
         "Is Collaborative Literacy 3rd Ed. + UFLI Foundations the right K-5 ELA choice "
         "for Troy's balanced-literacy → Science of Reading transition?",
         size=17, bold=True, color=GRAY_DARK)

# Answer "No." inline with header
add_text(s, Inches(0.5), Inches(2.25), Inches(2), Inches(0.32),
         "ANSWER", size=10, bold=True, color=ACCENT_RED)
add_text(s, Inches(0.5), Inches(2.55), Inches(3.0), Inches(1.0),
         "No.", size=64, bold=True, color=ACCENT_RED)
add_text(s, Inches(2.8), Inches(2.95), Inches(10), Inches(0.6),
         "The evidence converges from three independent directions:",
         size=15, bold=True, color=GRAY_DARK)

# Three convergent evidence columns
col_y = Inches(3.95)
col_w = Inches(4.05)
col_h = Inches(3.05)
col_x = [Inches(0.5), Inches(4.65), Inches(8.78)]

# Column 1: 0 of 50
add_rect(s, col_x[0], col_y, col_w, col_h, LIGHT_RED)
add_rect(s, col_x[0], col_y, col_w, Inches(0.06), ACCENT_RED)
add_text(s, col_x[0]+Inches(0.2), col_y+Inches(0.2), col_w-Inches(0.4), Inches(0.95),
         "0 of 50", size=52, bold=True, color=ACCENT_RED, align=PP_ALIGN.CENTER)
add_text(s, col_x[0]+Inches(0.2), col_y+Inches(1.25), col_w-Inches(0.4), Inches(0.8),
         "outperformer + peer districts¹ use Collaborative Literacy as their K-5 core.",
         size=12.5, bold=True, color=GRAY_DARK, align=PP_ALIGN.CENTER)
add_rect(s, col_x[0]+Inches(0.7), col_y+Inches(2.1), col_w-Inches(1.4), Inches(0.025), ACCENT_RED)
add_text(s, col_x[0]+Inches(0.2), col_y+Inches(2.2), col_w-Inches(0.4), Inches(0.8),
         "There is no national precedent for the combination Troy is considering.",
         size=11, color=GRAY_DARK, align=PP_ALIGN.CENTER, italic=True)

# Column 2: 19 of 20 DOTR outperformers
add_rect(s, col_x[1], col_y, col_w, col_h, LIGHT_GREEN)
add_rect(s, col_x[1], col_y, col_w, Inches(0.06), ACCENT_GREEN)
add_text(s, col_x[1]+Inches(0.2), col_y+Inches(0.2), col_w-Inches(0.4), Inches(0.95),
         "19 of 20", size=48, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
add_text(s, col_x[1]+Inches(0.2), col_y+Inches(1.25), col_w-Inches(0.4), Inches(0.8),
         "Education Scorecard 2026³ ELA-relevant outperformer districts run, supplement, or are transitioning to SoR-aligned curricula.",
         size=11.5, bold=True, color=GRAY_DARK, align=PP_ALIGN.CENTER)
add_rect(s, col_x[1]+Inches(0.7), col_y+Inches(2.1), col_w-Inches(1.4), Inches(0.025), ACCENT_GREEN)
add_text(s, col_x[1]+Inches(0.2), col_y+Inches(2.2), col_w-Inches(0.4), Inches(0.8),
         "Demographic-adjusted gain methodology — Harvard, Stanford, Dartmouth.\n\nOnly Wayne County NC stayed in pure balanced literacy.",
         size=10, color=GRAY_DARK, align=PP_ALIGN.CENTER)

# Column 3: 3 of 3 with named adoptions
add_rect(s, col_x[2], col_y, col_w, col_h, LIGHT_GREEN)
add_rect(s, col_x[2], col_y, col_w, Inches(0.06), ACCENT_GREEN)
add_text(s, col_x[2]+Inches(0.2), col_y+Inches(0.2), col_w-Inches(0.4), Inches(0.95),
         "3 of 3", size=52, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
add_text(s, col_x[2]+Inches(0.2), col_y+Inches(1.25), col_w-Inches(0.4), Inches(0.8),
         "of Troy's stronger alternatives have named district adoptions with recovery evidence.",
         size=11.5, bold=True, color=GRAY_DARK, align=PP_ALIGN.CENTER)
add_rect(s, col_x[2]+Inches(0.7), col_y+Inches(2.1), col_w-Inches(1.4), Inches(0.025), ACCENT_GREEN)
add_text(s, col_x[2]+Inches(0.2), col_y+Inches(2.2), col_w-Inches(0.4), Inches(0.8),
         "Amplify CKLA — Marion Co KY³⁰, Fond du Lac WI³³\nEL Education — Detroit DPSCD⁴²\nWit & Wisdom — W. Baton Rouge³⁷, Baltimore",
         size=10, color=GRAY_DARK, align=PP_ALIGN.CENTER, font="Consolas")

footer(s, 2)

# =================================================================
# SLIDE 3 — THE HONEST PRE-COVID PICTURE
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "First, the honest pre-COVID picture — Troy gained on aggregate, but subgroups were uneven",
          "SEDA cohort-standardized scale (cs): NAEP-anchored grade-level units. +1.0 = one grade above national norm.")

add_text(s, Inches(0.5), Inches(1.05), Inches(12), Inches(0.4),
         "Why this matters", size=12, bold=True, color=TROY_BLUE)
add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.8),
         ("Any argument for change has to start by acknowledging what's true: pre-COVID, Troy was healthy on the\n"
          "only metric that compares across state tests fairly¹⁻². Calkins UoS didn't break Troy historically.\n"
          "Pretending otherwise undermines the credibility of the post-COVID case (which IS strong)."),
         size=12, color=GRAY_DARK)

# Stats panel left
add_rect(s, Inches(0.5), Inches(2.6), Inches(6.0), Inches(4.2), LIGHT_GREEN)
add_text(s, Inches(0.7), Inches(2.75), Inches(5.6), Inches(0.4),
         "Troy SD 2009 → 2019 on SEDA cs scale²", size=13, bold=True, color=ACCENT_GREEN)
add_text(s, Inches(0.7), Inches(3.15), Inches(5.6), Inches(3.0),
         ("•  All Students:      +0.486 → +0.747   (+0.26)\n"
          "•  Asian:             +0.712 → +1.063   (+0.35)\n"
          "•  White:             +0.450 → +0.609   (+0.16)\n"
          "•  Black:             −0.070 → −0.160   (−0.09)\n"
          "•  Hispanic*:         +0.340 → +0.033   (−0.31)\n"
          "•  Econ Disadv:       −0.210 → +0.078   (+0.29)\n"
          "•  Not-Econ Disadv:   +0.552 → +0.853   (+0.30)\n\n"
          "* Hispanic suppressed pre-2011 (N<10); baseline is 2011."),
         size=11.5, color=GRAY_DARK, font="Consolas")
add_text(s, Inches(0.7), Inches(5.85), Inches(5.6), Inches(0.9),
         "Aggregate gains went to Asian, White, Not-ECD, and EconDis subgroups. "
         "Black + Hispanic subgroups stagnated or slipped slightly — the pre-COVID story was uneven.",
         size=11, color=GRAY_DARK, italic=True)

# Key takeaway right
add_rect(s, Inches(6.8), Inches(2.6), Inches(6.0), Inches(4.2), GRAY_LIGHT)
add_text(s, Inches(7.0), Inches(2.75), Inches(5.6), Inches(0.4),
         "What this rules out", size=13, bold=True, color=TROY_BLUE)
add_text(s, Inches(7.0), Inches(3.15), Inches(5.6), Inches(3.5),
         ("•  \"Calkins UoS has been failing Troy for a decade.\"\n"
          "    Mostly false. Aggregate + affluent gains 2009-2019.\n"
          "    (Caveat: Black + Hispanic subgroups stagnated.)\n\n"
          "•  \"This is a long-running trend that COVID exposed.\"\n"
          "    False. 2019 was Troy's peak on the aggregate scale.\n\n"
          "•  \"Troy was on a declining trajectory.\"\n"
          "    False. Troy outperformed every MI affluent peer\n"
          "    on aggregate 2009-2019.\n\n"
          "The case for change is NOT retrospective.\n"
          "The collapse is post-COVID (2019→2025) — and it hit\n"
          "every subgroup."),
         size=12, color=GRAY_DARK)

footer(s, 3)

# =================================================================
# SLIDE 4 — THE POST-COVID COLLAPSE
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "The post-COVID collapse is the real signal",
          "Troy's 2019 peak was +0.747. By 2025: +0.387. Lost 16 years of slow gains in 3 years.")
add_pic(s, f'{CHART_DIR}/chart10_troy_2009_2025.png',
        Inches(0.4), Inches(1.0), width=Inches(8.5))
# Right callouts — red box sized to actually contain ALL its text
add_rect(s, Inches(9.1), Inches(1.05), Inches(3.9), Inches(2.95), LIGHT_RED)
add_text(s, Inches(9.25), Inches(1.15), Inches(3.6), Inches(0.32),
         "From 2019 peak to 2025¹", size=12, bold=True, color=ACCENT_RED)
add_text(s, Inches(9.25), Inches(1.5), Inches(3.6), Inches(1.45),
         ("All Students    −0.36\n"
          "Asian           −0.46\n"
          "White           −0.29\n"
          "Not-ECD         −0.37\n"
          "EconDis         −0.25\n"
          "Black           −0.15\n"
          "Hispanic        +0.08 †"),
         size=12, color=GRAY_DARK, font="Consolas")
add_text(s, Inches(9.25), Inches(3.0), Inches(3.6), Inches(0.85),
         "(grade-level units)\n† small-N subgroup (~3-4% of Troy)",
         size=9.5, color=GRAY_MID, italic=True)

add_rect(s, Inches(9.1), Inches(4.1), Inches(3.9), Inches(2.65), GRAY_LIGHT)
add_text(s, Inches(9.25), Inches(4.2), Inches(3.6), Inches(0.35),
         "Four straight years of decline", size=12, bold=True, color=TROY_BLUE)
add_text(s, Inches(9.25), Inches(4.55), Inches(3.6), Inches(2.05),
         ("2022:  +0.605\n2023:  +0.526\n2024:  +0.390\n2025:  +0.387\n\n"
          "The 2025 Troy score (+0.387) is BELOW Troy's 2009 score (+0.486).\n\n"
          "Sixteen years of slow gains erased in three years — and recovery has not started."),
         size=11.5, color=GRAY_DARK)
footer(s, 4)

# =================================================================
# SLIDE 5 — NATIONAL RANKING: TROY 47 of 49
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "Of 50 demographically-comparable + outperformer districts, Troy ranks 47 of 49",
          "Pre-COVID (2017-19 mean) vs Post-COVID (2022-25 mean) on SEDA cohort-standardized RLA scale")
add_pic(s, f'{CHART_DIR}/chart11_seda_pre_post_covid.png',
        Inches(0.4), Inches(0.95), width=Inches(8.2))
# Right panel
add_rect(s, Inches(8.8), Inches(1.0), Inches(4.3), Inches(6.0), GRAY_LIGHT)
add_text(s, Inches(8.95), Inches(1.15), Inches(4.0), Inches(0.4),
         "Bottom 5 districts (Δ pre→post)", size=12, bold=True, color=ACCENT_RED)
add_text(s, Inches(8.95), Inches(1.55), Inches(4.0), Inches(1.4),
         ("48  Sikeston R-6         −0.275\n"
          "47  Detroit DPSCD        −0.259\n"
          "46  Troy SD              −0.253  ◄\n"
          "45  Bethlehem Area SD    −0.246\n"
          "44  West Bloomfield      −0.243"),
         size=10.5, color=GRAY_DARK, font="Consolas")
add_text(s, Inches(8.95), Inches(3.05), Inches(4.0), Inches(0.4),
         "Top 5 districts (Δ pre→post)", size=12, bold=True, color=ACCENT_GREEN)
add_text(s, Inches(8.95), Inches(3.45), Inches(4.0), Inches(1.4),
         ("1  Spring Branch ISD    +0.284  (TX STR³⁶)\n"
          "2  Palo Alto USD        +0.132  (ESRI³⁵)\n"
          "3  West Baton Rouge     +0.132  (W&W³⁷)\n"
          "4  Johnson City TN      +0.121  (TN HQIM³⁹)\n"
          "5  Frisco ISD           +0.097  (BL)"),
         size=10.5, color=GRAY_DARK, font="Consolas")
add_text(s, Inches(8.95), Inches(5.0), Inches(4.0), Inches(0.4),
         "Pattern", size=12, bold=True, color=TROY_BLUE)
add_text(s, Inches(8.95), Inches(5.4), Inches(4.0), Inches(1.8),
         ("Districts that GAINED post-COVID all run structured-literacy cores.\n\n"
          "Districts at the bottom either (a) face severe poverty headwinds (Sikeston, Detroit) or (b) are affluent BL districts (Troy, West Bloomfield)."),
         size=11.5, color=GRAY_DARK)
footer(s, 5)


# =================================================================
# SLIDE 6 (NEW) — COMPREHENSIVE CURRICULUM-TYPE BREAKDOWN
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "The comprehensive view — every district, every curriculum type",
          "Curriculum alone doesn't determine recovery. But magnitude of gain correlates with structured literacy.")

# Top summary box
add_rect(s, Inches(0.4), Inches(1.0), Inches(12.5), Inches(0.95), GRAY_LIGHT)
add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
         "Of 49 districts with sufficient pre/post-COVID data¹:", size=13, bold=True, color=TROY_BLUE)
add_text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.5),
         "13 gained on the SEDA scale  •  36 declined  •  Many SoR-adopters also declined — recovery is not automatic.",
         size=12, color=GRAY_DARK)

# Three-column breakdown
col_y = Inches(2.15)
col_h = Inches(4.55)
col_w = Inches(4.05)
col_x = [Inches(0.4), Inches(4.65), Inches(8.9)]

# Helper for slide-6 columns: header bar, big count, district rect-table, takeaway
def render_tier_column(cx, header_color, light_color, label, count_label,
                       districts, takeaway):
    add_rect(s, cx, col_y, col_w, col_h, light_color)
    add_rect(s, cx, col_y, col_w, Inches(0.45), header_color)
    add_text(s, cx+Inches(0.2), col_y+Inches(0.08), col_w-Inches(0.4), Inches(0.32),
             label, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, cx+Inches(0.2), col_y+Inches(0.55), col_w-Inches(0.4), Inches(0.35),
             count_label, size=14, bold=True, color=header_color, align=PP_ALIGN.CENTER)
    # district rect-table: 3 cols (name / Δ / tag)
    name_x = cx + Inches(0.2);  name_w = Inches(1.85)
    val_x  = cx + Inches(2.10); val_w  = Inches(0.85)
    tag_x  = cx + Inches(2.95); tag_w  = Inches(1.0)
    row_y = col_y + Inches(1.05)
    for name, val, tag in districts:
        add_text(s, name_x, row_y, name_w, Inches(0.24), name,
                 size=9.5, color=GRAY_DARK)
        add_text(s, val_x, row_y, val_w, Inches(0.24), val,
                 size=9.5, color=GRAY_DARK, font="Consolas")
        add_text(s, tag_x, row_y, tag_w, Inches(0.24), tag,
                 size=9, color=header_color, bold=True)
        row_y += Inches(0.26)
    # takeaway box at bottom of column
    take_top = col_y + col_h - Inches(1.15)
    add_text(s, cx+Inches(0.2), take_top, col_w-Inches(0.4), Inches(1.1),
             takeaway, size=9.5, color=GRAY_DARK)

# Big gain (≥+0.10)
render_tier_column(
    col_x[0], ACCENT_GREEN, LIGHT_GREEN,
    "Big gain   (Δ ≥ +0.10)",
    "4 districts — ALL SoR",
    [("Spring Branch ISD",  "+0.284", "TX STR³⁶"),
     ("Palo Alto USD",      "+0.132", "ESRI³⁵"),
     ("West Baton Rouge",   "+0.132", "W&W³⁷"),
     ("Johnson City TN",    "+0.121", "TN HQIM³⁹")],
    "100% of districts in this tier run structured-literacy programs.\n\n"
    "Next-strongest gainer outside this tier: Frisco ISD (TX, BL) at +0.097.")

# Moderate (+0.05 to +0.10)
render_tier_column(
    col_x[1], ACCENT_ORANGE, RGBColor(0xFF, 0xF4, 0xE0),
    "Moderate gain   (+0.05 to +0.10)",
    "6 districts — mixed",
    [("Frisco ISD",     "+0.097", "BL"),
     ("Milpitas USD",   "+0.090", "Benchmark"),
     ("Issaquah SD",    "+0.087", "Benchmark '24"),
     ("Dublin USD",     "+0.084", "Bench+Heggerty"),
     ("Coppell ISD",    "+0.068", "BL"),
     ("Princeton PS",   "+0.060", "BL (NJ peer)")],
    "Mix of BL, Benchmark-Advance, and recently-adopted structured literacy.")

# Flat or declined  (Δ < +0.05)
render_tier_column(
    col_x[2], ACCENT_RED, LIGHT_RED,
    "Flat or declined   (Δ < +0.05)",
    "39 districts — both types",
    [("Detroit DPSCD",   "−0.259", "EL Ed (SoR)"),
     ("Aldine ISD",      "−0.034", "CKLA (SoR)"),
     ("Baltimore PS",    "−0.213", "W&W (SoR)"),
     ("Fond du Lac",     "−0.014", "CKLA (SoR)"),
     ("Seaford SD",      "−0.180", "Bookworms (SoR)"),
     ("Troy SD",         "−0.253", "Calkins UoS (BL)")],
    "36 of 49 had Δ < 0; 3 held flat (Bloomfield Hills, Starkville, Dover). "
    "SoR alone is necessary but not sufficient — recovery also needs sustained PD, "
    "universal screeners, and a 2-3-year cohort lag (slide 19).")

# Bottom callout
add_text(s, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.4),
         "Honest finding: SoR doesn't guarantee recovery. But the only districts achieving STRONG recovery so far are SoR-aligned.",
         size=12.5, bold=True, color=TROY_BLUE, align=PP_ALIGN.CENTER)
footer(s, 6)

# =================================================================
# SLIDE 7 (NEW) — EDUCATION SCORECARD DOTR COMPREHENSIVE BREAKDOWN
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "A second independent benchmark — Education Scorecard 2026 'Districts on the Rise'",
          "Demographic-adjusted gain methodology (Harvard-Stanford-Dartmouth³) — catches outperformers SEDA raw Δ misses")

# Top context
add_rect(s, Inches(0.4), Inches(1.0), Inches(12.5), Inches(0.85), GRAY_LIGHT)
add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
         "What DOTR measures (different from SEDA raw Δ):", size=12, bold=True, color=TROY_BLUE)
add_text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.4),
         "Demographic-adjusted gain ≥ 0.3 grade levels in BOTH 2019–2025 AND 2022–2025, with ≥4 demographically-comparable peer districts in-state.  Detroit DPSCD is a DOTR outperformer despite a SEDA raw Δ decline — it gained more than its demographics predicted.",
         size=10.5, color=GRAY_DARK)

# Math + Reading group (17 districts) — real 3-column table for clean alignment
mr_left = Inches(0.4); mr_top = Inches(2.05); mr_w = Inches(8.2); mr_h = Inches(4.8)
add_rect(s, mr_left, mr_top, mr_w, mr_h, LIGHT_GREEN)
add_rect(s, mr_left, mr_top, mr_w, Inches(0.4), ACCENT_GREEN)
add_text(s, mr_left+Inches(0.15), mr_top+Inches(0.08), mr_w-Inches(0.3), Inches(0.3),
         "Math + Reading outperformers (17)", size=12.5, bold=True, color=WHITE)

# Column geometry
c1_x = mr_left + Inches(0.15)    # District
c1_w = Inches(2.0)
c2_x = c1_x + c1_w               # Curriculum
c2_w = Inches(4.85)
c3_x = c2_x + c2_w               # Type
c3_w = Inches(1.15)

mr_rows = [
    ("East Hartford CT",   "Fundations + Heggerty + Fund. Unl.", "SoR"),
    ("East Chicago IN",    "Heggerty + Reading Horizons + coaches", "SoR"),
    ("Marion County KY",   "Amplify CKLA + LETRS",               "SoR"),
    ("Baltimore City MD",  "Wit & Wisdom K-8 (2018-)",           "SoR"),
    ("Detroit DPSCD MI",   "EL Education K-5 (2018-)",           "SoR"),
    ("Starkville MS",      "MS HQIM (BL→structured shift)",      "SoR"),
    ("Pierre SD",          "Open Court + OG + Heggerty + Barton","SoR"),
    ("Johnson City TN",    "TN HQIM (SoR mandate)",              "SoR"),
    ("Spring Branch TX",   "TX STR framework",                   "SoR"),
    ("Fond du Lac WI",     "Amplify CKLA + DDI + UVA-PLE",       "SoR"),
    ("Dover NH",           "Knowledge-building shift",           "SoR-lean"),
    ("Brandywine DE",      "Explicit phonics K-2",               "SoR-lean"),
    ("Sikeston R-6 MO",    "Teacher-developed + state LETRS",    "SoR-lean"),
    ("Atlanta APS GA",     "Benchmark + HMH (Phase 2 SoR)",      "SoR-trans."),
    ("Kuna Joint ID",      "HMH Into Reading planned 2026-27",   "SoR-future"),
    ("College Comm. IA",   "MTSS focus (curriculum unclear)",    "Unclear"),
    ("Wayne County PS NC", "Stayed BL",                          "BL"),
]
row_h = Inches(0.235)
y = mr_top + Inches(0.5)
for district, curr, ctype in mr_rows:
    add_text(s, c1_x, y, c1_w, row_h, district, size=9, bold=True, color=TROY_BLUE)
    add_text(s, c2_x, y, c2_w, row_h, curr,     size=9, color=GRAY_DARK)
    color = ACCENT_GREEN if ctype == "SoR" else (ACCENT_ORANGE if "lean" in ctype or "trans" in ctype else (TROY_BLUE if ctype in ("SoR-future","Unclear") else ACCENT_RED))
    add_text(s, c3_x, y, c3_w, row_h, ctype, size=9, bold=True, color=color)
    y += row_h

# Reading-Only group (3)
add_rect(s, Inches(8.75), Inches(2.05), Inches(4.15), Inches(2.95), LIGHT_GREEN)
add_rect(s, Inches(8.75), Inches(2.05), Inches(4.15), Inches(0.4), ACCENT_GREEN)
add_text(s, Inches(8.9), Inches(2.13), Inches(4), Inches(0.3),
         "Reading-Only outperformers (3)", size=12.5, bold=True, color=WHITE)
add_text(s, Inches(8.9), Inches(2.55), Inches(4), Inches(2.4),
         ("Modesto City Elem CA\n"
          "  Benchmark Adv + UFLI + LETRS\n"
          "  → SoR-patched\n\n"
          "West Baton Rouge LA\n"
          "  Wit & Wisdom + Wilson Fundations\n"
          "  → SoR\n\n"
          "Roanoke County VA\n"
          "  W&W + Heggerty + Really Great Rdg\n"
          "  → SoR"),
         size=9.5, color=GRAY_DARK)

# Summary box
add_rect(s, Inches(8.75), Inches(5.15), Inches(4.15), Inches(1.65), LIGHT_RED)
add_text(s, Inches(8.9), Inches(5.22), Inches(4), Inches(0.32),
         "Summary across both groups", size=11, bold=True, color=ACCENT_RED)
add_text(s, Inches(8.9), Inches(5.55), Inches(4), Inches(0.95),
         ("• 13 districts:  clear SoR-aligned core\n"
          "• 5 districts:   SoR-leaning / patched / transitioning\n"
          "• 1 district:    future SoR adoption\n"
          "• 1 district:    unclear curriculum\n"
          "• 1 district:    stayed BL  (Wayne County NC)"),
         size=9, color=GRAY_DARK)
add_text(s, Inches(8.9), Inches(6.45), Inches(4), Inches(0.3),
         "19 of 20  =  95% SoR-aligned",
         size=10, bold=True, color=ACCENT_RED)

# Bottom callout
add_text(s, Inches(0.4), Inches(6.92), Inches(12.5), Inches(0.4),
         "Two independent methodologies converge: SEDA raw Δ + Education Scorecard demographic-adjusted gain both point at structured literacy.",
         size=11, bold=True, color=TROY_BLUE, align=PP_ALIGN.CENTER)
footer(s, 7)

# =================================================================
# SLIDE 8 — MI PEER COMPARISON: SUBGROUP ACTUAL VS EXPECTED
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "Of 8 Michigan affluent peers, Troy fell the furthest",
          "Same state, same M-STEP test, similar demographics. Pre-COVID expected vs post-COVID actual (SEDA cs Δ).")

# Left: line chart (taller aspect ratio, fills left side)
add_pic(s, f'{CHART_DIR}/chart12_mi_peers_2009_2025.png',
        Inches(0.1), Inches(1.0), width=Inches(8.1))

# Right panel: full subgroup Δ table with expected in parens
panel_x = Inches(8.3)
panel_w = Inches(4.85)
panel_top = Inches(1.0)
panel_h = Inches(6.1)
add_rect(s, panel_x, panel_top, panel_w, panel_h, GRAY_LIGHT)
add_text(s, panel_x + Inches(0.1), panel_top + Inches(0.05), Inches(4.6), Inches(0.3),
         "Pre→Post COVID Δ  (expected pre-COVID in parens)", size=9, bold=True, color=TROY_BLUE)

# Full subgroup data: (delta, pre_expected) per subgroup
# Subgroups: All, Asian, White, Black, Hispanic, EconDis, Not-ECD
subgroup_full = [
    ("Troy SD",         [(-0.253,0.730), (-0.347,1.103), (-0.231,0.592), (-0.274,-0.051), (-0.031,0.160), (-0.147,0.013), (-0.252,0.848)]),
    ("West Bloomfield", [(-0.243,0.167), (+0.068,0.433), (-0.125,0.381), (-0.274,-0.276), (-0.206,-0.085), (-0.249,-0.278), (-0.197,0.415)]),
    ("Rochester CSD",   [(-0.215,0.560), (-0.171,0.755), (-0.231,0.575), (-0.340,0.072), (-0.114,0.228), (-0.161,-0.081), (-0.208,0.637)]),
    ("Novi CSD",        [(-0.103,0.569), (-0.033,0.729), (-0.229,0.586), (-0.335,-0.083), (+0.164,0.158), (-0.236,-0.073), (-0.046,0.625)]),
    ("Northville PS",   [(-0.087,0.715), (-0.285,1.174), (-0.082,0.646), (None,None), (+0.226,0.291), (+0.024,-0.039), (-0.077,0.758)]),
    ("Birmingham PS",   [(-0.017,0.542), (-0.105,0.926), (-0.022,0.594), (-0.058,0.075), (-0.079,0.583), (+0.029,-0.021), (-0.008,0.592)]),
    ("Bloomfield Hills",[( 0.026,0.466), (-0.023,0.781), (+0.076,0.458), (-0.446,0.004), (None,None), (+0.014,-0.206), (+0.036,0.538)]),
]
sg_labels = ["All", "Asn", "Wht", "Blk", "Hsp", "EcD", "N-EC"]

# Top block: All, Asian, White, Black  (4 cols)
# Bottom block: Hispanic, EconDis, Not-ECD  (3 cols)
col_dist_x = panel_x + Inches(0.1)
col_dist_w = Inches(1.3)
row_h = Inches(0.3)

def render_block(start_y, sg_indices, sg_names):
    hdr_y = start_y
    add_rect(s, panel_x + Inches(0.05), hdr_y, panel_w - Inches(0.1), Inches(0.26), TROY_BLUE)
    add_text(s, col_dist_x, hdr_y, col_dist_w, Inches(0.26),
             "District", size=7.5, bold=True, color=WHITE)
    col_w = Inches((4.85 - 1.4) / len(sg_indices))
    col_xs = [panel_x + Inches(1.4) + col_w * i for i in range(len(sg_indices))]
    for i, idx in enumerate(sg_indices):
        add_text(s, col_xs[i], hdr_y, col_w, Inches(0.26),
                 sg_names[i], size=7.5, bold=True, color=WHITE, font="Consolas")
    ry = hdr_y + Inches(0.3)
    for dist_name, sg_vals in subgroup_full:
        is_troy = "Troy" in dist_name
        if is_troy:
            add_rect(s, panel_x + Inches(0.05), ry - Inches(0.01),
                     panel_w - Inches(0.1), row_h, LIGHT_RED)
        add_text(s, col_dist_x, ry, col_dist_w, row_h,
                 dist_name, size=7.5, bold=is_troy, color=ACCENT_RED if is_troy else GRAY_DARK)
        for i, idx in enumerate(sg_indices):
            delta, pre = sg_vals[idx]
            if delta is None:
                cell = "—"
                d_color = GRAY_MID
            else:
                d_color = ACCENT_RED if delta < -0.15 else (ACCENT_GREEN if delta > 0 else GRAY_DARK)
                cell = f"{delta:+.2f} ({pre:.2f})"
            add_text(s, col_xs[i], ry, col_w, row_h,
                     cell, size=7, font="Consolas",
                     color=d_color if is_troy or (delta is not None and abs(delta) > 0.15) else GRAY_DARK)
        ry += row_h
    return ry

# Block 1: All, Asian, White, Black
blk1_end = render_block(panel_top + Inches(0.4), [0, 1, 2, 3], ["All", "Asian", "White", "Black"])

# Block 2: Hispanic, EconDis, Not-ECD
render_block(blk1_end + Inches(0.15), [4, 5, 6], ["Hisp", "EconDis", "Not-ECD"])

footer(s, 8)


# =================================================================
# SLIDE 10 (NEW) — STUDENTS WITH DISABILITIES TRAJECTORY (M-STEP)
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "The Students-with-Disabilities subgroup tells the same story — only louder",
          "MI School Data M-STEP %Adv+Prof¹⁰ (different metric from SEDA cs scale — direct MDE source for Troy SWD)")

# Chart left — bigger, fills the vertical space
add_pic(s, f'{CHART_DIR}/chart_swd_deck.png',
        Inches(0.3), Inches(1.0), width=Inches(8.5))

# Right column: Troy SWD vs Peer-7 SWD delta table
add_rect(s, Inches(8.95), Inches(1.0), Inches(4.3), Inches(3.0), LIGHT_RED)
add_text(s, Inches(9.1), Inches(1.08), Inches(3.85), Inches(0.32),
         "Pre→Post-COVID Δ  (M-STEP %Adv+Prof)", size=11, bold=True, color=ACCENT_RED)
# 3-column table: Grade | Troy SWD | Peer-7 SWD
_swd_rows = [
    ("Grade 3", "−11.6 pp", "−3.8 pp", "−7.9 pp", True),
    ("Grade 4",  "−1.6 pp", "−0.1 pp", "−1.5 pp", False),
    ("Grade 5",  "−4.6 pp", "+0.7 pp", "−5.3 pp", True),
]
_cx_g = Inches(9.15); _cx_troy = Inches(10.20); _cx_peer = Inches(11.20); _cx_gap = Inches(12.15)
_cw = Inches(0.95)
add_text(s, _cx_g,    Inches(1.5), _cw, Inches(0.28), "Grade",    size=10, bold=True, color=GRAY_MID)
add_text(s, _cx_troy, Inches(1.5), _cw, Inches(0.28), "Troy",     size=10, bold=True, color=ACCENT_RED, align=PP_ALIGN.RIGHT)
add_text(s, _cx_peer, Inches(1.5), _cw, Inches(0.28), "Peer-7",   size=10, bold=True, color=TROY_BLUE, align=PP_ALIGN.RIGHT)
add_text(s, _cx_gap,  Inches(1.5), _cw, Inches(0.28), "Gap",      size=10, bold=True, color=GRAY_MID, align=PP_ALIGN.RIGHT)
add_rect(s, _cx_g, Inches(1.82), Inches(3.85), Inches(0.02), GRAY_MID)
_ry = Inches(1.95)
for gr, troy_v, peer_v, gap_v, hi in _swd_rows:
    color = ACCENT_RED if hi else GRAY_DARK
    add_text(s, _cx_g,    _ry, _cw, Inches(0.30), gr,     size=11, color=color, bold=hi)
    add_text(s, _cx_troy, _ry, _cw, Inches(0.30), troy_v, size=11, color=ACCENT_RED, bold=hi, font="Consolas", align=PP_ALIGN.RIGHT)
    add_text(s, _cx_peer, _ry, _cw, Inches(0.30), peer_v, size=11, color=TROY_BLUE, font="Consolas", align=PP_ALIGN.RIGHT)
    add_text(s, _cx_gap,  _ry, _cw, Inches(0.30), gap_v,  size=11, color=color, bold=hi, font="Consolas", align=PP_ALIGN.RIGHT)
    _ry += Inches(0.36)
add_text(s, Inches(9.1), Inches(3.05), Inches(3.85), Inches(0.6),
         "G3: Troy SWD fell 3× faster than peers.\nG5: Peers gained; Troy fell.",
         size=10, italic=True, color=ACCENT_RED)

# Why this matters
add_rect(s, Inches(8.95), Inches(3.75), Inches(4.3), Inches(2.9), GRAY_LIGHT)
add_text(s, Inches(9.1), Inches(3.82), Inches(3.85), Inches(0.32),
         "Why this matters", size=11, bold=True, color=TROY_BLUE)
add_text(s, Inches(9.1), Inches(4.15), Inches(3.85), Inches(2.4),
         ("•  Troy SWD G3: 37% pre-COVID → 25% today.\n"
          "   Peer-7 SWD G3: 35% → 31% (held much better).\n\n"
          "•  Pre-COVID Troy SWD ≈ Peer-7 SWD.\n"
          "   Post-COVID Troy SWD sits 4–7pp BELOW peers.\n\n"
          "•  SoR is engineered for struggling readers.\n"
          "   The students who benefit MOST from SoR\n"
          "   are falling FASTEST at Troy vs. same-state peers.\n\n"
          "•  Refutes \"regression-to-mean\" — SWD students\n"
          "   aren't high-achievers regressing."),
         size=9.5, color=GRAY_DARK)

# Bottom callout — full-width, styled like other slide callouts (not a thin band)
add_rect(s, Inches(0.4), Inches(6.78), Inches(12.5), Inches(0.30), TROY_BLUE)
add_text(s, Inches(0.4), Inches(6.83), Inches(12.5), Inches(0.25),
         "Two independent sources converge: SEDA cs scale (49 districts) and MDE M-STEP %prof (Troy subgroups) both say the SoR-relevant subgroups are falling fastest.",
         size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
footer(s, 10)

# =================================================================
# SLIDE 11 — THREE COMMON DEFENSES EMPIRICALLY REFUTED
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "The subgroup data refutes the three common defenses",
          "Each common defense of Troy's current trajectory is testable on the SEDA scale¹. Each fails the test.")
defenses = [
    ("1", "\"It's the ESL/ELL composition.\"",
     "If true: Hispanic and EL subgroups should be hit hardest.\n\n"
     "Data: Troy Hispanic Δ = −0.031. Rank 12 of 39 — best of Troy's subgroups¹.\n\n"
     "REFUTED. The Hispanic subgroup is the most resilient at Troy.", "#1F3A5F"),
    ("2", "\"It's the disadvantaged students.\"",
     "If true: EconDis subgroup should be hit harder than Not-ECD.\n\n"
     "Data: Troy EconDis Δ = −0.147 (rank 32/47). Troy Not-ECD Δ = −0.252 (rank 42/46)¹.\n\n"
     "REFUTED. Not-ECD declined MORE than EconDis.", "#CC6A11"),
    ("3", "\"It's a demographic / COVID story.\"",
     "If true: Affluent students with home-literacy environments should be PROTECTED.\n\n"
     "Data: Troy Asian Δ = −0.347 (34/35), White Δ = −0.231 (44/46) — the affluent subgroups dropped MOST¹.\n\n"
     "REFUTED. Affluence provides no protection at Troy. This points at instruction.", "#B02121"),
]
for i, (n, head, body, color) in enumerate(defenses):
    y = Inches(1.1) + i*Inches(2.0)
    add_rect(s, Inches(0.5), y, Inches(0.7), Inches(1.8), RGBColor.from_string(color[1:]))
    add_text(s, Inches(0.55), y+Inches(0.55), Inches(0.7), Inches(0.7),
             n, size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.4), y+Inches(0.05), Inches(11.5), Inches(0.5),
             head, size=15, bold=True, color=TROY_BLUE)
    add_text(s, Inches(1.4), y+Inches(0.55), Inches(11.5), Inches(1.3),
             body, size=13, color=GRAY_DARK)
footer(s, 11)

# =================================================================
# SLIDE 9 — ASIAN SUBGROUP DETAIL
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "Asian subgroup: Troy ranks 34 of 35 districts on post-COVID Δ",
          "Of 35 districts with Asian-subgroup data in SEDA 2025.1, only Bethlehem Area SD lost more ground.")
add_pic(s, f'{CHART_DIR}/chart17_asian_subgroup_rank.png',
        Inches(0.4), Inches(1.0), height=Inches(5.85))
# Right panel — shifted left to close gap with height-constrained chart
add_rect(s, Inches(7.0), Inches(1.0), Inches(6.0), Inches(5.85), GRAY_LIGHT)
add_text(s, Inches(7.15), Inches(1.15), Inches(5.7), Inches(0.4),
         "Asian subgroup top 5 (gainers)", size=12, bold=True, color=ACCENT_GREEN)
add_text(s, Inches(7.15), Inches(1.55), Inches(5.7), Inches(1.6),
         ("Spring Branch ISD    +0.717  (TX STR³⁶)\n"
          "Atlanta APS          +0.390  (SoR PD)\n"
          "Milpitas USD         +0.281\n"
          "Johnson City TN      +0.277\n"
          "Palo Alto USD        +0.269  (ESRI³⁵)"),
         size=11.5, color=GRAY_DARK, font="Consolas")
add_text(s, Inches(7.15), Inches(3.25), Inches(5.7), Inches(0.4),
         "Asian subgroup bottom 5", size=12, bold=True, color=ACCENT_RED)
add_text(s, Inches(7.15), Inches(3.65), Inches(5.7), Inches(1.6),
         ("Bethlehem Area SD    −0.426\n"
          "Troy SD              −0.347  ◄\n"
          "Northville PS        −0.285\n"
          "Wayne County PS      −0.264\n"
          "Modesto City Elem    −0.257"),
         size=11.5, color=GRAY_DARK, font="Consolas")
add_text(s, Inches(7.15), Inches(5.4), Inches(5.7), Inches(0.4),
         "Direct affluent-peer signal", size=12, bold=True, color=TROY_BLUE)
add_text(s, Inches(7.15), Inches(5.8), Inches(5.7), Inches(1.0),
         ("Palo Alto (Troy's closest demographic match by Asian %) exited Calkins UoS in 2021 via ESRI³⁵.\n\n"
          "Result: +0.269 grade-level gain on Asian subgroup. Direct evidence the shift works."),
         size=11.5, color=GRAY_DARK)
footer(s, 12)

# =================================================================
# SLIDE 10 — NOT-ECD DETAIL
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "Not-Econ-Disadvantaged subgroup: Troy ranks 42 of 46",
          "Troy's affluent students dropped −0.252 grade levels post-COVID — more than 90% of districts in our universe.")
add_pic(s, f'{CHART_DIR}/chart16_troy_notecd_trend.png',
        Inches(0.4), Inches(1.0), width=Inches(8.5))
# Right
add_rect(s, Inches(9.1), Inches(1.0), Inches(3.9), Inches(6.0), GRAY_LIGHT)
add_text(s, Inches(9.25), Inches(1.15), Inches(3.6), Inches(0.4),
         "Not-ECD top 5 (recovery)", size=11, bold=True, color=ACCENT_GREEN)
add_text(s, Inches(9.25), Inches(1.55), Inches(3.6), Inches(1.6),
         ("West Baton Rouge    +0.315\n"
          "Spring Branch       +0.193\n"
          "Palo Alto USD       +0.138\n"
          "Frisco ISD          +0.133\n"
          "Issaquah SD         +0.124"),
         size=11.5, color=GRAY_DARK, font="Consolas")
add_text(s, Inches(9.25), Inches(3.2), Inches(3.6), Inches(0.4),
         "Not-ECD bottom 5", size=11, bold=True, color=ACCENT_RED)
add_text(s, Inches(9.25), Inches(3.6), Inches(3.6), Inches(1.6),
         ("Bethlehem Area      −0.326\n"
          "Wayne County PS     −0.302\n"
          "Seaford SD          −0.301\n"
          "Brandywine SD       −0.267\n"
          "Troy SD             −0.252  ◄"),
         size=11.5, color=GRAY_DARK, font="Consolas")
add_text(s, Inches(9.25), Inches(5.3), Inches(3.6), Inches(0.4),
         "Why this matters", size=11, bold=True, color=TROY_BLUE)
add_text(s, Inches(9.25), Inches(5.7), Inches(3.6), Inches(1.4),
         ("These are the families with the most home-literacy resources. If Troy's curriculum design were sound, these are the kids who should recover first. They're recovering last."),
         size=10, color=GRAY_DARK)
footer(s, 13)

# =================================================================
# SLIDE 11 (NEW) — CKLA DISTRICTS LEAPFROGGING THEIR STATES
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "CKLA-adopting districts on the MDE-approved list are leapfrogging their states",
          "Marion County KY and Fond du Lac WI both adopted Amplify CKLA. Both gained vs their state averages while Troy lost ground.")
add_pic(s, f'{CHART_DIR}/chart19_ckla_cohort_effect.png',
        Inches(0.4), Inches(0.95), width=Inches(8.5))
# Right panel
add_rect(s, Inches(9.1), Inches(1.0), Inches(3.9), Inches(2.95), LIGHT_GREEN)
add_text(s, Inches(9.25), Inches(1.15), Inches(3.6), Inches(0.4),
         "State-relative gain after CKLA¹", size=12, bold=True, color=ACCENT_GREEN)
add_text(s, Inches(9.25), Inches(1.55), Inches(3.6), Inches(2.4),
         ("Marion County KY:\n"
          "  Pre-CKLA vs state    +0.02\n"
          "  Post-CKLA vs state   +0.19\n"
          "  Δ vs state           +0.17\n\n"
          "Fond du Lac WI:\n"
          "  Pre-CKLA vs state    −0.02\n"
          "  Post-CKLA vs state   +0.10\n"
          "  Δ vs state           +0.12"),
         size=12, color=GRAY_DARK, font="Consolas")
add_rect(s, Inches(9.1), Inches(4.1), Inches(3.9), Inches(2.95), LIGHT_RED)
add_text(s, Inches(9.25), Inches(4.25), Inches(3.6), Inches(0.4),
         "Troy on the same scale", size=12, bold=True, color=ACCENT_RED)
add_text(s, Inches(9.25), Inches(4.65), Inches(3.6), Inches(2.4),
         ("Troy SD MI (Calkins UoS):\n"
          "  Pre-COVID vs state   +0.79\n"
          "  2025 vs state        +0.67\n"
          "  Δ vs state           −0.11\n\n"
          "Troy is losing ~0.11 grade levels relative to MI state.\n"
          "The two CKLA peers gained +0.12 and +0.17 relative to their states."),
         size=11, color=GRAY_DARK, font="Consolas")
footer(s, 14)

# =================================================================
# SLIDE 12 (NEW) — MARION COUNTY + FOND DU LAC PLAYBOOKS
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "How they did it — Marion County KY and Fond du Lac WI playbooks",
          "Both districts adopted Amplify CKLA (Tier-1 on the MI MDE Section 35m approved list¹⁸)")

# Marion County KY column
add_rect(s, Inches(0.5), Inches(1.0), Inches(6.1), Inches(5.9), LIGHT_GREEN)
add_text(s, Inches(0.65), Inches(1.1), Inches(5.8), Inches(0.4),
         "Marion County KY — Chris Brady, Supt.³⁰,³²", size=14, bold=True, color=ACCENT_GREEN)
add_text(s, Inches(0.65), Inches(1.5), Inches(5.8), Inches(0.32),
         "~2,937 students  •  61% econ-disadvantaged  •  83% White, 11% Hispanic",
         size=10, color=GRAY_DARK, italic=True)

mc_rows = [
    ("Curriculum stack",
     "Amplify CKLA K-5 + Amplify ELA 6-8 (single-vendor vertical alignment). District-wide rollout in 2022-23, aligned with KY Read to Succeed Act of 2022."),
    ("Teacher PD",
     "LETRS Volumes 1 & 2 via Kentucky Reading Academies³¹ (state-paid, no cost to district) — 44 teachers AND administrators trained since 2022. Principals went through LETRS WITH teachers so they can evaluate foundational instruction in the same vocabulary."),
    ("Coaching restructure",
     "Coaches moved from building-based to content-driven, operating in 6-9 WEEK CYCLES with teachers on standards-aligned instruction and diagnostic data. The single most unusual structural choice — most rural KY districts kept building-based coaching."),
    ("Implementation",
     "Full district-wide rollout, NOT a pilot. Tech use reduced to reclaim student-teacher engagement time. G3-7 ELA proficient +10pp 2021-22 → 2024-25."),
    ("Result on SEDA scale",
     "Pre-CKLA: −0.01 vs KY state. 2024-25: +0.19 vs state. +0.17 grade-level state-relative gain in 3 years. Marion County MS jumped 221→57 in KY state rank."),
]
for i, (head, body) in enumerate(mc_rows):
    y = Inches(2.0) + i*Inches(0.96)
    add_text(s, Inches(0.65), y, Inches(5.8), Inches(0.3),
             head, size=12.5, bold=True, color=TROY_BLUE)
    add_text(s, Inches(0.65), y+Inches(0.32), Inches(5.8), Inches(0.65),
             body, size=11, color=GRAY_DARK)

# Fond du Lac WI column
add_rect(s, Inches(6.8), Inches(1.0), Inches(6.1), Inches(5.9), LIGHT_GREEN)
add_text(s, Inches(6.95), Inches(1.1), Inches(5.8), Inches(0.4),
         "Fond du Lac WI — Matt Steinbarth, Supt.³³,³⁴", size=14, bold=True, color=ACCENT_GREEN)
add_text(s, Inches(6.95), Inches(1.5), Inches(5.8), Inches(0.32),
         "~6,457 students  •  49.5% econ-disadvantaged  •  59% White, 21% Hispanic, 9% Black",
         size=10, color=GRAY_DARK, italic=True)

fdl_rows = [
    ("Curriculum stack",
     "Amplify CKLA K-5 + Amplify ELA 6-8 + Bridges Math + TCI Social Studies + Mystery Science (K-8 coherent stack). Adoption aligned with August 2022 Strategic Plan."),
    ("Leadership development partnership",
     "Partnered with the UVA Darden Partnership for Leaders in Education (UVA-PLE)³⁴. Single biggest differentiator: CKLA embedded in a multi-year leadership and instructional-redesign program, NOT adopted as a standalone literacy textbook."),
    ("Teacher PD",
     "Weekly Data-Driven Instruction (DDI) cycles trained via UVA-PLE. ~$20M ESSER funds deployed for curriculum + after-school + classroom libraries + AVID. WI Act 20 (2023) added a state legal forcing function³⁸."),
    ("Result",
     "Jumped from 375th to 175th of 421 WI districts. \"Exceeds Expectations\" on 2024-25 WI report card — first time since 2016-17. Highest reading growth of all 421 WI districts 2022-2024."),
    ("Leadership transition",
     "Supt Fleig led the redesign 2022-2025; retired May 2025. Board promoted CAO Matt Steinbarth (25-yr district insider) 7-0 to protect continuity. Credit the turnaround to Fleig + Steinbarth + UVA-PLE."),
]
for i, (head, body) in enumerate(fdl_rows):
    y = Inches(2.0) + i*Inches(0.96)
    add_text(s, Inches(6.95), y, Inches(5.8), Inches(0.3),
             head, size=12.5, bold=True, color=TROY_BLUE)
    add_text(s, Inches(6.95), y+Inches(0.32), Inches(5.8), Inches(0.65),
             body, size=11, color=GRAY_DARK)

footer(s, 15)

# =================================================================
# SLIDE 13 — RECOVERY CASE STUDIES
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "Recovery is possible — and it tracks with curriculum choice",
          "Four districts gained post-COVID. All four run structured-literacy programs.")
add_pic(s, f'{CHART_DIR}/chart14_subgroup_compare.png',
        Inches(0.4), Inches(0.95), width=Inches(8.5))
# Right panel
cards = [
    ("Spring Branch ISD (TX)³⁶", "+0.284 overall / +0.717 Asian / +0.193 Not-ECD",
     "TX STR-aligned framework, K paraprofessionals trained in evidence-based instruction."),
    ("Palo Alto USD (CA)³⁵", "+0.132 overall / +0.269 Asian / +0.138 Not-ECD",
     "Exited Calkins Units of Study in 2021 via Every Student Reads Initiative. Direct Troy parallel."),
    ("West Baton Rouge (LA)³⁷", "+0.132 overall / +0.315 Not-ECD / +0.242 White / +0.079 Black",
     "Wit & Wisdom + Wilson Fundations K-3. K-3 literacy 57% → 68% in single year."),
    ("Johnson City TN³⁹", "+0.121 overall / +0.277 Asian / +0.165 White",
     "Tennessee HQIM SoR mandate. +22pp above TN state on G4 ELA in 2025."),
]
for i, (head, deltas, body) in enumerate(cards):
    y = Inches(1.0) + i*Inches(1.55)
    add_rect(s, Inches(9.1), y, Inches(3.9), Inches(1.45), LIGHT_GREEN)
    add_text(s, Inches(9.25), y+Inches(0.05), Inches(3.6), Inches(0.3),
             head, size=12, bold=True, color=ACCENT_GREEN)
    add_text(s, Inches(9.25), y+Inches(0.35), Inches(3.6), Inches(0.3),
             deltas, size=10, bold=True, color=GRAY_DARK, font="Consolas")
    add_text(s, Inches(9.25), y+Inches(0.7), Inches(3.6), Inches(0.7),
             body, size=11, color=GRAY_DARK)
footer(s, 16)

# =================================================================
# SLIDE 12 — LONG BEACH NATURAL EXPERIMENT (POST-COVID)
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "Long Beach USD — the natural experiment is now empirically visible",
          "Same district, same Carl Cohn-era leadership culture, same demographics. Only the curriculum changed.")
add_pic(s, f'{CHART_DIR}/chart13_post_covid_recovery.png',
        Inches(0.4), Inches(0.95), width=Inches(8.5))
# Right
add_rect(s, Inches(9.1), Inches(1.0), Inches(3.9), Inches(3.0), GRAY_LIGHT)
add_text(s, Inches(9.25), Inches(1.15), Inches(3.6), Inches(0.4),
         "Long Beach trajectory", size=12, bold=True, color=TROY_BLUE)
add_text(s, Inches(9.25), Inches(1.55), Inches(3.6), Inches(2.3),
         ("1992-2002 (Cohn era):  Open Court basal\n"
          "→ Broad Prize 2003\n\n"
          "2010s drift:  Teachers College Units of Study (Calkins)\n"
          "→ same program Troy adopted 2014\n\n"
          "Pre/Post-COVID Δ:  −0.076 grade levels"),
         size=11.5, color=GRAY_DARK)
add_rect(s, Inches(9.1), Inches(4.15), Inches(3.9), Inches(2.85), LIGHT_RED)
add_text(s, Inches(9.25), Inches(4.3), Inches(3.6), Inches(0.4),
         "Why this matters for Troy", size=12, bold=True, color=ACCENT_RED)
add_text(s, Inches(9.25), Inches(4.7), Inches(3.6), Inches(2.2),
         ("Long Beach is the closest thing to a controlled trial against Calkins UoS in U.S. K-5 ELA data.\n\n"
          "Same district. Same insider supt tradition. Same demographics. The pre-COVID period worked.\n\n"
          "Post-COVID, the Calkins-derived core failed to support recovery — same pattern as Troy."),
         size=11.5, color=GRAY_DARK)
footer(s, 17)

# =================================================================
# SLIDE 13 — THE MECHANISM
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "Why balanced literacy fails when school and home reading routines both break down",
          "A plain-English causal story that's consistent with every subgroup finding")

# Left column: How Calkins is designed
add_rect(s, Inches(0.5), Inches(1.1), Inches(4.0), Inches(5.7), GRAY_LIGHT)
add_text(s, Inches(0.65), Inches(1.25), Inches(3.7), Inches(0.4),
         "How Calkins UoS works²⁰,²³", size=13, bold=True, color=TROY_BLUE)
add_text(s, Inches(0.65), Inches(1.65), Inches(3.7), Inches(5.0),
         ("•  Workshop-model, student-led discovery\n\n"
          "•  Leveled-text independent reading\n\n"
          "•  Comprehension strategies (predict, infer, etc.)\n\n"
          "•  Minimal explicit phonics in core\n\n"
          "•  Relies on rich home language environment to fill gaps in explicit instruction\n\n"
          "•  Reading-aloud, book access, parent-supported reading time"),
         size=12.5, color=GRAY_DARK)

# Middle column: What broke
add_rect(s, Inches(4.7), Inches(1.1), Inches(4.0), Inches(5.7), LIGHT_RED)
add_text(s, Inches(4.85), Inches(1.25), Inches(3.7), Inches(0.4),
         "What COVID broke²¹", size=13, bold=True, color=ACCENT_RED)
add_text(s, Inches(4.85), Inches(1.65), Inches(3.7), Inches(5.0),
         ("•  Classroom delivery — interrupted\n\n"
          "•  Home reading routines — disrupted\n\n"
          "•  Book access at home — limited\n\n"
          "•  Parent time for reading — fragmented\n\n"
          "•  Balanced literacy assumed BOTH school AND home would teach reading — when home routines broke down, school's part alone wasn't designed to do the whole job\n\n"
          "•  Most-affluent students had the MOST home reading support to lose, so they fell furthest"),
         size=12.5, color=GRAY_DARK)

# Right column: Why SoR was protected
add_rect(s, Inches(8.9), Inches(1.1), Inches(4.0), Inches(5.7), LIGHT_GREEN)
add_text(s, Inches(9.05), Inches(1.25), Inches(3.7), Inches(0.4),
         "Why SoR districts recovered²²,²⁵", size=13, bold=True, color=ACCENT_GREEN)
add_text(s, Inches(9.05), Inches(1.65), Inches(3.7), Inches(5.0),
         ("•  Explicit phonics scope and sequence\n\n"
          "•  Decodable texts at K-2\n\n"
          "•  Knowledge-building core that doesn't assume background knowledge\n\n"
          "•  Cohesive sequence — same content in every classroom\n\n"
          "•  Doesn't depend on home support — designed to teach reading in school\n\n"
          "•  When classroom delivery returned, the curriculum was still intact and could rebuild skills"),
         size=12.5, color=GRAY_DARK)
footer(s, 18)

# =================================================================
# SLIDE 18 (NEW) — GRADE-LEVEL Δ COUNTER-HYPOTHESIS TEST
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "Tested hypothesis — does SoR help only early grades, not later ones?",
          "If the 'BL for later grades' hypothesis were correct, BL districts should hold up at G5. They don't.")

# Top: the hypothesis being tested
add_rect(s, Inches(0.4), Inches(1.0), Inches(12.5), Inches(0.7), GRAY_LIGHT)
add_text(s, Inches(0.6), Inches(1.05), Inches(2.0), Inches(0.3),
         "The hypothesis", size=11.5, bold=True, color=TROY_BLUE)
add_text(s, Inches(0.6), Inches(1.32), Inches(12.2), Inches(0.35),
         "\"Maybe SoR helps for K-2 phonics, but balanced literacy still has a role at G3-G5 where comprehension + rich text matter more.\" — pedagogically plausible. Testable on this data.",
         size=10.5, color=GRAY_DARK, italic=True)

# Left: mean Δ by grade and curriculum type
add_rect(s, Inches(0.4), Inches(1.85), Inches(6.2), Inches(4.6), LIGHT_GREEN)
add_text(s, Inches(0.55), Inches(1.95), Inches(6), Inches(0.4),
         "Mean SEDA Δ by grade × curriculum type", size=12.5, bold=True, color=ACCENT_GREEN)
add_text(s, Inches(0.55), Inches(2.4), Inches(6), Inches(4.0),
         ("Type              G3 Δ      G4 Δ      G5 Δ      n\n"
          "─────────────────────────────────────────────────\n"
          "SoR-clear         −0.088   −0.117   −0.065    21\n"
          "SoR-recent/lean   +0.031   −0.011   −0.036     5\n"
          "Mixed/Benchmark   +0.113   +0.133   +0.062     3\n"
          "Balanced lit.     −0.102   −0.107   −0.083    15\n\n"
          "Reading:\n"
          "• BL districts decline at EVERY grade — they don't hold up\n"
          "  better at G5. They decline less, but still decline.\n"
          "• SoR-clear shows its SMALLEST decline at G5, not G3.\n"
          "• SoR-recent adopters (Bellevue, Issaquah, Lake Wash.,\n"
          "  Modesto, College Comm.) actually GAINED at G3 (+0.031)\n"
          "  — cohort effect from new-K students reaching G3 in '24-25."),
         size=10, color=GRAY_DARK, font="Consolas")

# Right: Top 5 G5 gainers
add_rect(s, Inches(6.75), Inches(1.85), Inches(6.15), Inches(4.6), LIGHT_GREEN)
add_text(s, Inches(6.9), Inches(1.95), Inches(6), Inches(0.4),
         "Top 5 G5 gainers — 100% SoR-aligned", size=12.5, bold=True, color=ACCENT_GREEN)
add_text(s, Inches(6.9), Inches(2.4), Inches(6), Inches(4.0),
         ("1  Spring Branch ISD    +0.284   TX STR\n"
          "2  Johnson City TN      +0.221   TN HQIM\n"
          "3  West Baton Rouge     +0.207   W&W + Fundations\n"
          "4  Steubenville City    +0.176   Success for All (25 yr)\n"
          "5  Starkville-Oktibbeha +0.126   MS HQIM\n\n"
          "Top 5 G3 gainers (for comparison):\n"
          "1  Palo Alto USD        +0.205   SoR (ESRI)\n"
          "2  West Baton Rouge     +0.142   SoR\n"
          "3  Milpitas USD         +0.099   Mixed/Benchmark\n"
          "4  Issaquah SD          +0.087   SoR-recent\n"
          "5  College Community    +0.075   SoR-recent\n\n"
          "The strongest G5 gainers are EXCLUSIVELY SoR.\n"
          "The strongest G3 gainers are 4/5 SoR + 1 mixed."),
         size=10, color=GRAY_DARK, font="Consolas")

# Bottom: caveats
add_rect(s, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.42), LIGHT_RED)
add_text(s, Inches(0.55), Inches(6.60), Inches(2.8), Inches(0.3),
         "Confounders this data can't resolve:", size=10, bold=True, color=ACCENT_RED)
add_text(s, Inches(3.4), Inches(6.60), Inches(9.5), Inches(0.35),
         "Cohort effects (G5 had K 5 yrs ago); cross-grade scale differences; COVID hit early grades hardest (K-1 missed foundational year); small SoR-recent/Mixed n.",
         size=9, color=GRAY_DARK)
footer(s, 19)

# =================================================================
# SLIDE 19 — CURRICULUM VERDICT MATRIX
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "Curriculum verdict — Troy's options on the MI Section 35m list",
          "Ranked by evidence strength + alignment with recovery patterns")

table_left = Inches(0.4)
table_top = Inches(1.05)
col_widths = [Inches(0.6), Inches(3.8), Inches(2.3), Inches(2.8), Inches(2.3), Inches(0.8)]
headers = ["#", "Curriculum", "EdReports", "Recovery proof points", "Comp.", "Grade"]
add_rect(s, table_left, table_top, sum(col_widths, Emu(0)), Inches(0.4), TROY_BLUE)
x = table_left
for i, h in enumerate(headers):
    add_text(s, x, table_top+Inches(0.05), col_widths[i], Inches(0.3),
             h, size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER if i!=1 else PP_ALIGN.LEFT)
    x += col_widths[i]
rows = [
    ("1","Amplify CKLA 3rd Ed.","All Green¹⁴","Aldine⁴¹ G4+11pp; Fond du Lac #1 WI³³","Single","A+",ACCENT_GREEN),
    ("2","EL Education K-5","All Green¹⁴","Detroit DPSCD⁴²; NYC Reads Dist 11⁴³","Single","A",ACCENT_GREEN),
    ("3","Wit & Wisdom + Fundations + Geodes","All Green¹⁴","West Baton Rouge³⁷; Baltimore W&W since 2018","2-3 vendor","A",ACCENT_GREEN),
    ("4","UFLI Foundations (component)","Not reviewed¹⁴","RRQ 2025 ES >1.0 SD²⁴","Supplement","A",ACCENT_GREEN),
    ("5","Really Great Reading","196/198¹⁴","NYC Reads foundational core⁴³","Supplement","A−",ACCENT_GREEN),
    ("…","(…)","","","","",GRAY_MID),
    ("11","Collaborative Literacy 3rd Ed. + UFLI","2016 ed.: Partial¹⁴","Zero of 50 outperformers use it¹","2-vendor","B / B−",ACCENT_RED),
    ("12","HMH Into Reading 2025","All Green (v1.0)¹⁴","Contested by APM Reports²¹","Single","B−",GRAY_MID),
]
row_h = Inches(0.42)
for ri, row in enumerate(rows):
    y = table_top + Inches(0.4) + ri*row_h
    highlight = 'Collab' in row[1]
    if highlight:
        add_rect(s, table_left, y, sum(col_widths, Emu(0)), row_h, LIGHT_RED)
    else:
        add_rect(s, table_left, y, sum(col_widths, Emu(0)), row_h, WHITE if ri%2==0 else GRAY_LIGHT)
    x = table_left
    for ci, val in enumerate(row[:6]):
        if ci == 5:
            color = row[6]; bold = True
        elif ci == 0:
            color = TROY_BLUE; bold = True
        else:
            color = GRAY_DARK; bold = highlight
        add_text(s, x, y+Inches(0.1), col_widths[ci], Inches(0.3),
                 val, size=11.5, bold=bold, color=color,
                 align=PP_ALIGN.CENTER if ci!=1 else PP_ALIGN.LEFT)
        x += col_widths[ci]

add_text(s, Inches(0.4), Inches(5.4), Inches(12.5), Inches(0.4),
         "Collab Lit + UFLI is the weakest comprehensive option on the Section 35m list for a district moving OUT of balanced literacy.",
         size=14, bold=True, color=ACCENT_RED)
add_text(s, Inches(0.4), Inches(5.85), Inches(12.5), Inches(1.4),
         ("Of 50 districts in the post-COVID national benchmark, zero use Collaborative Literacy as their K-5 core.\n"
          "Of the 4 districts that GAINED on the SEDA scale 2017-19 → 2022-25, all four use structured literacy.\n"
          "Of Troy's stronger alternatives — CKLA, EL Education, Wit & Wisdom — each has direct post-COVID recovery evidence."),
         size=12.5, color=GRAY_DARK)
footer(s, 20)

# =================================================================
# SLIDE 15 — RECOMMENDATION
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "Recommendation — three paths ranked",
          "Each leg of the evidence base points the same direction")

paths = [
    ("Recommended primary path", "Amplify CKLA 3rd Edition K-5",
     "Single vendor, integrated foundations + knowledge.  All-green EdReports (first ELA program to earn it).\n"
     "Marion County KY³⁰: +10pp G3-7 ELA 2021-25 with CKLA + LETRS.  Fond du Lac WI³³: #1 reading growth statewide.",
     ACCENT_GREEN),
    ("Strong alternative — best of breed", "Wit & Wisdom + Wilson Fundations K-3 + Geodes",
     "West Baton Rouge LA³⁷: K-3 literacy 57%→68% in ONE year. Best Not-ECD recovery (+0.315) in our universe.\n"
     "Roanoke County VA³: closest demographic peer to Troy. Baltimore City: longest urban track record (since 2018).",
     ACCENT_GREEN),
    ("Honorable mention", "EL Education K-5 (+ phonics supplement)",
     "Direct in-state Michigan validator: Detroit DPSCD⁴² adopted 2018, outpaced state.\n"
     "NYC Reads natural experiment⁴³: District 11 (EL Education) outperformed both Into Reading and Wit & Wisdom.",
     TROY_BLUE),
    ("If teacher buy-in or local politics blocks a clean change", "The least-bad fallback",
     "Keep Collab Lit as a transitional core; drop Making Meaning in G3-5; substitute a knowledge-building supplement.\n"
     "Use SIPPS for Tier-2 only and UFLI as Tier-1 K-2 phonics (avoid duplication). Commit to a 2-3 year sunset review.",
     GRAY_MID),
]
for i, (label, name, body, color) in enumerate(paths):
    y = Inches(1.05) + i*Inches(1.45)
    add_rect(s, Inches(0.5), y, Inches(0.15), Inches(1.35), color)
    add_text(s, Inches(0.8), y+Inches(0.05), Inches(3.5), Inches(0.3),
             label.upper(), size=10, bold=True, color=color)
    add_text(s, Inches(0.8), y+Inches(0.32), Inches(11.5), Inches(0.4),
             name, size=16, bold=True, color=TROY_BLUE)
    add_text(s, Inches(0.8), y+Inches(0.75), Inches(11.5), Inches(0.65),
             body, size=12, color=GRAY_DARK)
footer(s, 21)

# =================================================================
# SLIDE 16 — EXECUTION + TRACK
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "Curriculum is necessary but not sufficient — and what to track",
          "Five execution factors + the real-time peer benchmarks")

# Left: 5 execution factors
add_text(s, Inches(0.5), Inches(1.0), Inches(7), Inches(0.4),
         "Five execution factors in every winning district", size=14, bold=True, color=TROY_BLUE)
factors = [
    ("1", "Stack, don't bolt", "Knowledge core + systematic phonics + teacher PD as ONE system"),
    ("2", "60+ hours structured PD", "LETRS most-cited³¹ (3 of 20 DOTRs); brand varies; depth doesn't"),
    ("3", "Universal screeners + MTSS", "MDE-approved K-3 screeners: mCLASS/DIBELS 8, Amira, MAP Reading Fluency"),
    ("4", "Plan for a 2-3 year lag", "Bethlehem K-change 2015 → G3 bounce 2018²⁶"),
    ("5", "Fidelity is the binding constraint", "NYC Phase 1 conf 55% vs Phase 2 38%⁴³ = real outcome gap"),
]
for i, (n, head, body) in enumerate(factors):
    y = Inches(1.45) + i*Inches(0.85)
    add_rect(s, Inches(0.5), y, Inches(0.45), Inches(0.75), TROY_BLUE)
    add_text(s, Inches(0.5), y+Inches(0.15), Inches(0.45), Inches(0.5),
             n, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.1), y, Inches(6.0), Inches(0.35),
             head, size=13.5, bold=True, color=TROY_BLUE)
    add_text(s, Inches(1.1), y+Inches(0.36), Inches(6.0), Inches(0.4),
             body, size=11.5, color=GRAY_DARK)

# Right: districts to track
add_text(s, Inches(7.8), Inches(1.0), Inches(5.3), Inches(0.4),
         "Track these districts in real time", size=14, bold=True, color=TROY_BLUE)
tracks = [
    ("Lake Washington WA", "Board adoption Spring 2026", "Closest demographic peer mid-shift"),
    ("Bellevue SD WA", "Year 2 of ARC Core + UFLI + Heggerty", "Year-2 SBA fall 2026"),
    ("Issaquah SD WA", "First post-Benchmark Advance SBA", "Fall 2026"),
    ("Palo Alto USD CA", "Year 5 post-ESRI CAASPP trend", "Annual CAASPP"),
    ("Fond du Lac WI", "Continued CKLA + DDI trajectory", "WI Forward fall 2026"),
    ("Aldine ISD TX", "Watch G5+G6 STAAR 2026 (full post-CKLA cohort)", "TEA TAPR fall 2026"),
]
for i, (name, status, when) in enumerate(tracks):
    y = Inches(1.45) + i*Inches(0.85)
    add_rect(s, Inches(7.8), y, Inches(0.15), Inches(0.75), ACCENT_GREEN)
    add_text(s, Inches(8.0), y, Inches(5.0), Inches(0.3),
             name, size=12.5, bold=True, color=TROY_BLUE)
    add_text(s, Inches(8.0), y+Inches(0.3), Inches(5.0), Inches(0.25),
             status, size=11, color=GRAY_DARK)
    add_text(s, Inches(8.0), y+Inches(0.55), Inches(5.0), Inches(0.25),
             "→ " + when, size=10, color=ACCENT_GREEN, italic=True)
footer(s, 22)


# =================================================================
# SLIDE 19 — APPENDIX TITLE + METHODOLOGY
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "Appendix — Methodology",
          "Datasets, metrics, time windows, district universe")

# Two columns: Datasets used + Methodology choices
add_text(s, Inches(0.5), Inches(1.05), Inches(6), Inches(0.4),
         "Primary metric", size=14, bold=True, color=TROY_BLUE)
add_text(s, Inches(0.5), Inches(1.45), Inches(6), Inches(1.6),
         ("SEDA cohort-standardized score (cs): NAEP-anchored grade-level units. "
          "0.0 = at-grade national norm; +1.0 = one grade above national norm.\n\n"
          "This is the only metric that is genuinely cross-state comparable in absolute terms — it normalizes M-STEP, CAASPP, STAAR, LEAP, MAAP, TCAP, NJSLA, SBA, OST against a common NAEP-anchored scale."),
         size=12, color=GRAY_DARK)

add_text(s, Inches(0.5), Inches(3.3), Inches(6), Inches(0.4),
         "Time windows", size=14, bold=True, color=TROY_BLUE)
add_text(s, Inches(0.5), Inches(3.7), Inches(6), Inches(1.6),
         ("Pre-COVID baseline:  2017-2019 (3-year mean)\n"
          "Post-COVID window:  2022-2025 (4-year mean)\n"
          "Δ = post − pre (units: grade levels)\n\n"
          "2020 and 2021 omitted from SEDA — state tests were canceled or partially administered nationwide. This is the SEDA team's standard recovery-window convention."),
         size=12, color=GRAY_DARK, font="Consolas")

add_text(s, Inches(0.5), Inches(5.5), Inches(6), Inches(0.4),
         "Grade pooling", size=14, bold=True, color=TROY_BLUE)
add_text(s, Inches(0.5), Inches(5.9), Inches(6), Inches(1.0),
         ("G3-G5 ELA only. Pooled within district × year, weighted by tested-student count per grade. Math excluded (this analysis is about ELA)."),
         size=12, color=GRAY_DARK)

# Right column - District universe
add_text(s, Inches(7), Inches(1.05), Inches(6), Inches(0.4),
         "District universe — 50 targeted, 49 with valid SEDA Δ", size=13, bold=True, color=TROY_BLUE)
add_text(s, Inches(7), Inches(1.45), Inches(6), Inches(5.0),
         ("•  Troy SD + 7 MI affluent peers    (8)\n   (Bloomfield Hills, Birmingham, Northville, Novi, Rochester, Walled Lake, West Bloomfield)\n"
          "•  4 CA workbook peers + 4 CA outperformers    (8)\n   (Palo Alto, Milpitas, Walnut Valley, Dublin + Modesto, Sanger, Garden Grove, Long Beach)\n"
          "•  6 TX districts    (6)\n   (Aldine, Brownsville, Coppell, Plano, Frisco, Spring Branch)\n"
          "•  3 NJ peers + 3 WA SoR-shift peers    (6)\n   (WW-Plainsboro, Millburn, Princeton + Bellevue, Issaquah, Lake Washington)\n"
          "•  Steubenville (OH gold standard) + Detroit DPSCD (MI in-state)    (2)\n"
          "•  15 additional Education Scorecard 2026 DOTR districts    (15)\n   (CT, DE, GA, IA, ID, IN, KY, MD, MO, MS, NC, NH, SD, TN, WI — Detroit + Spring Branch already counted)\n"
          "•  2 DOTR Reading-Only adds    (2)   (West Baton Rouge LA, Roanoke County VA — Modesto already in CA)\n"
          "•  3 sustained / early-SoR comparators    (3)   (Seaford DE, Valley Stream 30 NY, Bethlehem PA)"),
         size=10.5, color=GRAY_DARK)
add_text(s, Inches(7), Inches(6.4), Inches(6), Inches(0.5),
         "Counts referenced elsewhere in the deck:", size=10.5, bold=True, color=TROY_BLUE)
add_text(s, Inches(7), Inches(6.65), Inches(6), Inches(0.5),
         "20 = ELA-relevant DOTR districts (17 Math+Reading + 3 Reading-only).  28 = predecessor state-test analysis (raw % proficient on M-STEP/CAASPP/STAAR/etc., before SEDA expansion) — see reports/quantitative_analysis.md.",
         size=9.5, color=GRAY_DARK, italic=True)

footer(s, 23)


# =================================================================
# SLIDE 23b — SEDA ELA PEER RANKING (APPENDIX)
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "Appendix — Troy vs. 504 level-matched peers, K-5 ELA (SEDA 2025.1)",
          "Pre-COVID baseline ±0.25 of Troy, enrollment ≥150/grade-year, complete 2017-2025 data")

add_pic(s, f'{CHART_DIR}/chart_seda_ela_peers.png',
        Inches(0.3), Inches(0.95), width=Inches(8.3))

# Right panel: key stats
add_text(s, Inches(8.8), Inches(1.0), Inches(4.3), Inches(0.4),
         "Key findings", size=14, bold=True, color=TROY_BLUE)

findings = [
    ("504", "level-matched peer districts nationwide"),
    ("81 / 504", "Troy's rank (16th percentile) — bottom sixth"),
    ("-0.252", "Troy's post-COVID Δ (grade levels vs national norm)"),
    ("-0.107", "Peer median Δ — Troy underperforms by 0.145 grade levels"),
]
for i, (num, desc) in enumerate(findings):
    y = Inches(1.5) + i * Inches(0.75)
    add_text(s, Inches(8.8), y, Inches(1.2), Inches(0.35),
             num, size=18, bold=True, color=ACCENT_RED)
    add_text(s, Inches(10.0), y, Inches(3.1), Inches(0.65),
             desc, size=11, color=GRAY_DARK)

# Subgroup breakdown
add_text(s, Inches(8.8), Inches(4.6), Inches(4.3), Inches(0.35),
         "Subgroup ranks (vs peers with data)", size=12, bold=True, color=TROY_BLUE)
sub_rows = [
    ("All Students", "81 / 504", "16th pct"),
    ("Asian", "21 / 257", "8th pct"),
    ("White", "96 / 496", "19th pct"),
    ("Econ Disadvantaged", "229 / 429", "53rd pct"),
]
for i, (sub, rank, pct) in enumerate(sub_rows):
    y = Inches(5.05) + i * Inches(0.4)
    color = ACCENT_RED if "8th" in pct or "16th" in pct or "19th" in pct else ACCENT_GREEN
    add_text(s, Inches(8.8), y, Inches(2.0), Inches(0.35),
             sub, size=11, color=GRAY_DARK)
    add_text(s, Inches(10.8), y, Inches(1.3), Inches(0.35),
             rank, size=11, bold=True, color=color)
    add_text(s, Inches(12.1), y, Inches(1.0), Inches(0.35),
             pct, size=10, color=color, italic=True)

add_text(s, Inches(8.8), Inches(6.7), Inches(4.3), Inches(0.5),
         "Peers derived independently from the Math analysis (44% overlap).\n"
         "Troy's ELA baseline (+0.729) is lower than Math (+0.937) — different peer universe.",
         size=9, color=GRAY_MID, italic=True)

footer(s, 24)

# =================================================================
# SLIDE 20 — PRIMARY DATASETS [refs 1-13]
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "References (1 of 5) — Primary datasets",
          "Each cited claim in the deck is anchored to one of these sources")

refs1 = [
    ("1", "Stanford Education Data Archive 2025.1", "Reardon, S. F., et al. (2026). SEDA v2025.1. https://edopportunity.org/trends/data/downloads/  •  Coverage 2009-2025; cohort-standardized scale anchored to NAEP."),
    ("2", "SEDA 6.0 (geographic district file)", "Reardon, S. F., et al. (2026). SEDA v6.0. https://edopportunity.org/opportunity/data/downloads/  •  Coverage 2009-2019; used to validate pre-COVID baseline."),
    ("3", "Education Scorecard 2026 \"Districts on the Rise\"", "Harvard / Stanford / Dartmouth. https://educationscorecard.org/districts-on-the-rise/  •  17 demographically-adjusted post-COVID outperformer districts."),
    ("4", "Troy achievement workbook (8 districts)", "tsd-achievement.karpowitsch.org master workbook  •  G3-G7 × subgroup × year, 2019-2025, 5,497 demographic cells."),
    ("5", "CAASPP Smarter Balanced Research Files", "https://caaspp-elpac.ets.org/caaspp/ResearchFileListSB  •  CA district-level ELA, 2019-2025; used for Modesto, Sanger, Garden Grove, Long Beach."),
    ("6", "Texas TAPR / TEA", "https://rptsvr1.tea.texas.gov/perfreport/tapr/  •  STAAR reading G3-G5 for Aldine, Brownsville, Spring Branch, Plano, Frisco, Coppell."),
    ("7", "Louisiana DOE LEAP", "https://doe.louisiana.gov/students/assessments/  •  West Baton Rouge Parish 2022-2025 LEAP Mastery+."),
    ("8", "Mississippi DOE MAAP", "https://reports.mdek12.org/  •  Starkville-Oktibbeha 2019-2025 (via tpcref.org compilation)."),
    ("9", "Tennessee Department of Education TCAP", "https://www.tn.gov/education/  •  Johnson City Schools 2019-2025 district assessment files."),
    ("10", "Michigan School Data CEPI / M-STEP", "https://www.mischooldata.org/  •  Troy SD trajectory + statewide G3-G5 averages."),
    ("11", "Ohio Department of Education / OST", "education.ohio.gov  •  Steubenville City Schools + statewide G3-G5 OST averages."),
    ("12", "Washington OSPI Smarter Balanced", "https://washingtonstatereportcard.ospi.k12.wa.us/  •  Bellevue / Issaquah / Lake Washington SBA."),
    ("13", "Oakland County 115 — MI affluent peer 2025", "https://oaklandcounty115.com/2025/09/14/test-results-show-reading-proficiency-among-third-and-fourth-graders-across-oakland-county/  •  G3+G4 ELA 2025 head-to-head."),
]
y = Inches(1.0)
for n, title, body in refs1:
    add_rect(s, Inches(0.4), y, Inches(0.55), Inches(0.42), TROY_BLUE)
    add_text(s, Inches(0.4), y+Inches(0.06), Inches(0.55), Inches(0.3),
             n, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.1), y, Inches(11.7), Inches(0.22),
             title, size=11.5, bold=True, color=TROY_BLUE)
    add_text(s, Inches(1.1), y+Inches(0.22), Inches(11.7), Inches(0.22),
             body, size=9, color=GRAY_DARK)
    y += Inches(0.45)
footer(s, 25)

# =================================================================
# SLIDE 21 — CURRICULUM EVIDENCE [refs 14-19]
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "References (2 of 5) — Curriculum evidence sources",
          "EdReports / Reading League / ESSA / WWC / MDE / Knowledge Matters")

refs2 = [
    ("14", "EdReports.org curriculum reviews", "https://edreports.org/  •  Green/Yellow/Red gateway ratings: Text Quality / Building Knowledge / Foundational Skills / Usability. v2.0 rubric (2024-25) explicitly assesses SoR alignment."),
    ("15", "The Reading League — Curriculum Navigation Reports", "https://www.thereadingleague.org/curriculum-navigation-reports/  •  SoR-only evaluation: \"red flag practices\" (three-cueing, leveled-text decoding, lack of phoneme-grapheme systematicity)."),
    ("16", "Evidence for ESSA", "https://www.evidenceforessa.org/  •  ESSA tiers (Strong, Moderate, Promising, Rationale) for K-5 ELA programs. Used for 95 Phonics (ESSA Strong), IMSE OG+ (Promising), etc."),
    ("17", "What Works Clearinghouse / IES", "https://ies.ed.gov/ncee/wwc/  •  Federal evidence registry. Reviewed Orton-Gillingham-based interventions; found insufficient qualifying studies."),
    ("18", "Michigan MDE Section 35m Tier-1 list", "https://www.michigan.gov/mde/services/academic-standards/literacy/literacy-grants/section-35m  •  State-approved K-5 evidence-based ELA curricula (Dec 2025 release)."),
    ("19", "Knowledge Matters Campaign", "https://knowledgematterscampaign.org/  •  Recognized knowledge-building K-5 ELA programs: CKLA, EL Education, Wit & Wisdom, ARC Core, Bookworms, Fishtank."),
]
y = Inches(1.0)
for n, title, body in refs2:
    add_rect(s, Inches(0.4), y, Inches(0.55), Inches(0.7), TROY_BLUE)
    add_text(s, Inches(0.4), y+Inches(0.2), Inches(0.55), Inches(0.3),
             n, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.1), y+Inches(0.05), Inches(11.7), Inches(0.3),
             title, size=13, bold=True, color=TROY_BLUE)
    add_text(s, Inches(1.1), y+Inches(0.32), Inches(11.7), Inches(0.4),
             body, size=10.5, color=GRAY_DARK)
    y += Inches(0.8)
footer(s, 26)

# =================================================================
# SLIDE 22 — RESEARCH & JOURNALISM [refs 20-29]
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "References (3 of 5) — Science of Reading research & journalism",
          "Independent scholarly + investigative sources cited in mechanism slide and curriculum evaluation")

refs3 = [
    ("20", "Emily Hanford / APM Reports", "\"Sold a Story\" podcast series + investigations. https://features.apmreports.org/sold-a-story/  •  Documented Calkins UoS / F&P / Reading Recovery critiques."),
    ("21", "APM Reports — EdReports & SoR (March 2025)", "https://www.apmreports.org/story/2025/03/06/edreports-reading-curriculum-reviews-science-of-reading  •  Investigation of EdReports all-green ratings on basal programs (Wonders, myView, Into Reading)."),
    ("22", "Natalie Wexler — The Knowledge Gap", "Avery (2019). https://nataliewexler.com/the-knowledge-gap/  •  Framework that comprehension depends on background knowledge built via coherent content sequence."),
    ("23", "Tim Shanahan — Shanahan on Literacy", "https://www.shanahanonliteracy.com/  •  Reading research and instructional practice. Strategy instruction critique cited in mechanism slide."),
    ("24", "Mark Seidenberg — Language at the Speed of Sight", "Basic Books (2017). https://seidenbergreading.net/  •  Cognitive science of reading; critique of balanced literacy."),
    ("25", "Lane et al. (2025) — UFLI Foundations RRQ study", "Reading Research Quarterly. https://ila.onlinelibrary.wiley.com/doi/10.1002/rrq.607  •  ES > 1.0 SD on early-literacy outcomes; districtwide Florida pilot."),
    ("26", "Karin Chenoweth — Districts That Succeed", "Harvard Education Press (2021)  •  Sustained-outperformer district case studies (Steubenville, Seaford, Lane OK, etc.)."),
    ("27", "Success for All Foundation — Steubenville case study", "https://www.successforall.org/wp-content/uploads/2025/05/CS_Steubenville-1.pdf  •  25-year SFA implementation; 93-100% G3 reading proficient."),
    ("28", "The 74 Million — Steubenville coverage", "https://www.the74million.org/article/why-steubenville-ohio-might-be-the-best-school-district-in-america/  •  Independent journalism on Steubenville."),
    ("29", "Karen Vaites — Southern Surge analysis", "https://www.karenvaites.org/p/the-southern-surge-understanding  •  Mississippi / Louisiana / Tennessee NAEP gains tied to SoR adoption."),
]
y = Inches(1.0)
for n, title, body in refs3:
    add_rect(s, Inches(0.4), y, Inches(0.55), Inches(0.55), TROY_BLUE)
    add_text(s, Inches(0.4), y+Inches(0.12), Inches(0.55), Inches(0.3),
             n, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.1), y, Inches(11.7), Inches(0.25),
             title, size=12, bold=True, color=TROY_BLUE)
    add_text(s, Inches(1.1), y+Inches(0.25), Inches(11.7), Inches(0.32),
             body, size=10, color=GRAY_DARK)
    y += Inches(0.6)
footer(s, 27)

# =================================================================
# SLIDE 23 — DISTRICT CASE STUDIES [refs 30-43]
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "References (4 of 5) — District case-study sources",
          "Specific districts cited in the deck: adoption + outcomes")

refs4 = [
    ("30", "Kentucky Teacher — Marion County coverage (May 2026)", "https://www.kentuckyteacher.org/news/2026/05/kentucky-ranks-among-nations-best-for-student-recovery-5th-in-reading-8th-in-mathematics/"),
    ("31", "Kentucky Reading Academies / Lexia LETRS", "https://www.lexialearning.com/kentucky-letrs  •  https://www.education.ky.gov/curriculum/EarlyLiteracy/Pages/ky_reading_Academies.aspx"),
    ("32", "Education Scorecard PDFs — case studies", "Marion County: https://educationscorecard.org/wp-content/uploads/2026/05/KY_Marion-County_Case-Study.pdf  •  20 DOTR district PDFs in /research/escorecard_pdfs/"),
    ("33", "Fond du Lac SD — Strategic Plan + Studer Education", "https://www.fonddulac.k12.wi.us/page/strategic-plan  •  Aug 2022 strategic plan; CKLA + Amplify ELA + Bridges + TCI + Mystery Science stack."),
    ("34", "UVA-PLE / Darden Partnership for Leaders in Education", "https://www.darden.virginia.edu/uva-ple/  •  Leadership and instructional-redesign partnership underlying Fond du Lac's turnaround."),
    ("35", "Palo Alto USD — Every Student Reads Initiative (ESRI)", "https://www.pausd.org/learning/esri  •  Exited Calkins UoS 2021; mandatory Orton-Gillingham K-3 training."),
    ("36", "Spring Branch ISD — Texas STR framework", "https://www.springbranchisd.com/studentsfamilies/literacy  •  TX Science of Teaching Reading state mandate."),
    ("37", "West Baton Rouge / LA DOE Wit & Wisdom adoption", "https://doe.louisiana.gov/docs/default-source/curricular-resources/great-minds-pbc---wit-wisdom-with-fundations-and-geodes-ela-grades-k-3-(-2023).pdf"),
    ("38", "WI Act 20 / Wisconsin Reads", "https://dpi.wi.gov/wi-reads/curriculum  •  State science-of-reading mandate (2023) with approved curriculum list."),
    ("39", "Tennessee Department of Education — Reading 360 + HQIM mandate", "https://www.tn.gov/education/reading-360.html  •  State SoR PD (60 hours) + approved K-5 curriculum list."),
    ("40", "EdWeek — \"Inside the Long Beach Way\" (2007)", "https://www.edweek.org/leadership/inside-the-long-beach-way/2007/09  •  Carl Cohn era Open Court tradition documented."),
    ("41", "Amplify — Aldine ISD CKLA adoption (Jul 2020)", "https://amplify.com/news/aldine-isd-adopts-amplifys-integrated-early-literacy-suite-supports-students-in-learning-to-read-confidently-by-third-grade/"),
    ("42", "Detroit DPSCD + EL Education — Skillman Foundation grant", "https://www.skillman.org/blog/standardswork-receives-grant-to-support-implementation-of-transformative-literacy-curriculum-in-detroit-public-schools/  •  https://www.detroitk12.org/academics/core-curriculum/english-language-arts-and-literacy"),
    ("43", "NYC Reads — Chalkbeat / EdWeek / K-12 Dive", "Chalkbeat: https://www.chalkbeat.org/newyork/2025/07/07/nyc-reads-solves-screener-data-eric-adams-school-curriculum-mandate/  •  EdWeek: Phase 1 vs Phase 2 confidence data (Jul 2025)."),
]
y = Inches(1.0)
for n, title, body in refs4:
    add_rect(s, Inches(0.4), y, Inches(0.55), Inches(0.42), TROY_BLUE)
    add_text(s, Inches(0.4), y+Inches(0.06), Inches(0.55), Inches(0.3),
             n, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.1), y, Inches(11.7), Inches(0.22),
             title, size=11, bold=True, color=TROY_BLUE)
    add_text(s, Inches(1.1), y+Inches(0.22), Inches(11.7), Inches(0.22),
             body, size=8.5, color=GRAY_DARK)
    y += Inches(0.45)
footer(s, 28)

# =================================================================
# SLIDE 24 — PROJECT ARTIFACTS + FILE MAP
# =================================================================
s = prs.slides.add_slide(BLANK)
title_bar(s, "References (5 of 5) — Project artifacts + how to verify",
          "Every claim is reproducible from the public GitHub repository")

add_text(s, Inches(0.5), Inches(0.95), Inches(12.5), Inches(0.4),
         "github.com/akarpo/tsd-k5ela-choice", size=13, bold=True, color=TROY_BLUE, font="Consolas")

# Two-column file map (paths are relative to the repo root)
left_groups = [
    ("Web presentation", [
        ("README.md", "Project README + slide map"),
        ("index.html", "Vanilla-JS slide viewer"),
        ("slides/*.png", "28 rendered slide PNGs (144 DPI)"),
        ("deck/Troy_K5_ELA_Executive_Summary.pptx", "Source PowerPoint deck"),
    ]),
    ("Reproducibility pipeline", [
        ("analysis/extract_seda_subset.py", "SEDA 2025.1 → 50-district subset"),
        ("analysis/compute_deltas.py", "Pre/post-COVID Δ matrix"),
        ("analysis/build_deck.py", "Builds the .pptx deck"),
        ("docs/REPRODUCIBILITY.md", "How to reproduce in 3 commands"),
        ("docs/METHODOLOGY.md", "Metrics, time windows, district universe"),
        ("docs/DISTRICT_UNIVERSE.md", "All 50 districts with NCES LEA IDs"),
    ]),
]
right_groups = [
    ("Processed data (data/)", [
        ("seda_2025_pooled.json", "SEDA cs G3-G5 pooled by year × subgroup"),
        ("seda_2025_state.json", "State-level SEDA averages by year"),
        ("seda_subgroup_delta.csv", "Pre/post-COVID Δ per subgroup × district"),
        ("master_dataset.csv", "2,544 rows: dist × year × grade × subgroup"),
        ("seda_2009_2025_extract.csv", "Raw SEDA 2025.1 extract (50 dists × G3-G5)"),
        ("troy_swd_ela.csv", "Troy SWD M-STEP %prof, G3-G5 2015-25"),
    ]),
    ("Supporting reports (reports/)", [
        ("synthesis.md", "Master synthesis report"),
        ("quantitative_analysis.md", "28-district state-test analysis"),
        ("education_scorecard_2026_dotr.md", "20 DOTR outperformer case studies"),
        ("curriculum_evaluations.md", "14 Section 35m curricula evaluated"),
        ("collab_lit_ufli_pressure_test.md", "Collab Lit + UFLI critique"),
    ]),
]

def render_column(col_groups, x_left):
    y = Inches(1.55)
    for group, files in col_groups:
        add_text(s, x_left, y, Inches(6.0), Inches(0.3),
                 group, size=13, bold=True, color=ACCENT_GREEN)
        y += Inches(0.35)
        for fname, desc in files:
            add_text(s, x_left+Inches(0.15), y, Inches(3.6), Inches(0.22),
                     fname, size=8.5, color=GRAY_DARK, font="Consolas", bold=True)
            add_text(s, x_left+Inches(3.8), y, Inches(2.5), Inches(0.22),
                     desc, size=8.5, color=GRAY_MID)
            y += Inches(0.26)
        y += Inches(0.15)

add_rect(s, Inches(0.4), Inches(1.4), Inches(6.4), Inches(5.5), GRAY_LIGHT)
add_rect(s, Inches(6.95), Inches(1.4), Inches(6.0), Inches(5.5), GRAY_LIGHT)
render_column(left_groups, Inches(0.55))
render_column(right_groups, Inches(7.10))

footer(s, 29)

# Save
import os
_HERE = os.path.dirname(os.path.abspath(__file__))
out = os.path.join(_HERE, "Troy_K5_ELA_Executive_Summary.pptx")
prs.save(out)
print(f"Wrote: {out}")
print(f"Slides: {len(prs.slides)}")
